from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation(filename="MoveX_Project_Presentation.pptx"):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "MoveX Cab Application"
    subtitle.text = "Complete Workflow, Features, and Technologies\nDetailed Project Report"

    # Slide 1: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Introduction & Overview"
    tf = body_shape.text_frame
    tf.text = "MoveX is a scalable and intuitive ride-hailing platform."
    p = tf.add_paragraph()
    p.text = "Comprises three primary components:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Customer App: For booking, tracking, and wallet management."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Driver App: For accepting rides, earnings, and navigation."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Admin Dashboard: For operational control and verification."
    p.level = 2

    # Slide 2: Technologies Used
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Technologies Used"
    tf = body_shape.text_frame
    tf.text = "Frontend / Mobile App:"
    p = tf.add_paragraph()
    p.text = "React Native, Expo, React Navigation, Context API"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend & APIs:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Node.js, Express.js"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database & Real-time:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "MongoDB, Mongoose (Geospatial queries)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Socket.io (Live tracking & Ride dispatching)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Authentication & Hosting:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Firebase Auth, Render (Cloud deployment)"
    p.level = 1

    # Slide 3: Customer App Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Customer App Features"
    tf = body_shape.text_frame
    tf.text = "Interactive Ride Booking with precise map selections."
    p = tf.add_paragraph()
    p.text = "Real-time Fare Estimation across vehicle tiers (Bike, Auto, Mini, Sedan, SUV)."
    p = tf.add_paragraph()
    p.text = "Subscription Passes (Tiered system for discounts and free rides)."
    p = tf.add_paragraph()
    p.text = "Live Tracking via Socket.io."
    p = tf.add_paragraph()
    p.text = "In-App Digital Wallet for cashless transactions."

    # Slide 4: Driver App & Admin Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Driver & Admin Features"
    tf = body_shape.text_frame
    tf.text = "Driver App:"
    p = tf.add_paragraph()
    p.text = "Toggle Online/Offline status."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Receive sequential ride requests with countdown timers."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Secure Document Upload (KYC, RC, Insurance)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Admin Dashboard:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Verify and approve driver profiles."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Manage dynamic pricing and pass configurations."
    p.level = 1

    # Slide 5: Complete Booking Workflow
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Complete Booking Workflow"
    tf = body_shape.text_frame
    tf.text = "1. Customer requests a ride (Geospatial query finds drivers)."
    p = tf.add_paragraph()
    p.text = "2. Backend sequences available drivers based on proximity."
    p = tf.add_paragraph()
    p.text = "3. Socket.io dispatches 'ride:incoming' event sequentially."
    p = tf.add_paragraph()
    p.text = "4. Driver accepts and navigates to pickup."
    p = tf.add_paragraph()
    p.text = "5. OTP Verification starts the trip."
    p = tf.add_paragraph()
    p.text = "6. Real-time location streaming during the trip."
    p = tf.add_paragraph()
    p.text = "7. Payment processing and user ratings complete the cycle."

    # Save Presentation
    prs.save(filename)
    print(f"Successfully generated {filename}")

if __name__ == "__main__":
    create_presentation()
